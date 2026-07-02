from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
DARK1 = RGBColor(0x07, 0x0B, 0x15)
DARK2 = RGBColor(0x0D, 0x14, 0x28)
DARK3 = RGBColor(0x11, 0x1B, 0x2E)
DARK4 = RGBColor(0x16, 0x22, 0x40)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
TEXT2 = RGBColor(0x94, 0xA3, 0xB8)
TEXT3 = RGBColor(0x47, 0x56, 0x69)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

def set_bg(slide, color=DARK1):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, color, left=0, top=0, width=None, height=None, radius=None):
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_txt(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, name='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb

def add_card(slide, title, desc, left, top, width, height, color=DARK3, title_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color or ACCENT
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT2
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    return shape

def add_accent_bar(slide, left, top, width=1.0, height=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

# ========================
# SLIDE 1: Title
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, RGBColor(0x05, 0x09, 0x12))
add_rect(slide, ACCENT, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

add_txt(slide, "✦  FINAL PROJECT PRESENTATION", 1.5, 1.4, 10, 0.5, 13, TEXT2, False, PP_ALIGN.CENTER)
add_txt(slide, "WAREHOUSE OS", 1.5, 2.0, 10, 1.5, 56, WHITE, True, PP_ALIGN.CENTER)
add_accent_bar(slide, 5.8, 3.6, 1.6)
add_txt(slide, "Multi-Agent Order Fulfillment System", 1.5, 3.9, 10, 0.7, 28, TEXT2, False, PP_ALIGN.CENTER)
add_txt(slide, "AI-Powered Warehouse Management with Autonomous Agents", 1.5, 4.5, 10, 0.5, 18, TEXT3, False, PP_ALIGN.CENTER)

tags = ["Python 3.12 + FastAPI", "Next.js 16 + React 19", "OpenAI Agents SDK", "PostgreSQL + Qdrant"]
for i, t in enumerate(tags):
    x = 2.5 + i * 2.5
    add_rect(slide, DARK3, Inches(x), Inches(5.3), Inches(2.2), Inches(0.45), 0.15)
    add_txt(slide, t, x + 0.1, 5.33, 2.0, 0.4, 12, TEXT2, False, PP_ALIGN.CENTER)

# ========================
# SLIDE 2: Project Overview
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, DARK2, Inches(0), Inches(0), Inches(5.8), Inches(7.5))

add_txt(slide, "✦  OVERVIEW", 0.6, 0.6, 5, 0.4, 13, ACCENT, False)
add_txt(slide, "Project Overview", 0.6, 1.0, 5, 0.7, 30, WHITE, True)
add_accent_bar(slide, 0.6, 1.7, 1.0)
add_txt(slide, "An intelligent warehouse management system powered by AI agents that automates the entire order fulfillment lifecycle — from order placement to delivery tracking.", 0.6, 2.1, 4.8, 1.8, 14, TEXT2)

features = [
    "Real-time order processing & tracking",
    "Multi-agent AI orchestration",
    "Vector database for semantic search",
    "Automated rerouting & delay prediction",
    "SMS/Email notifications (Twilio + SendGrid)",
    "3D warehouse visualization (Three.js)",
    "Smooth animations (Framer Motion + GSAP)",
]
tb = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(4.8), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
for i, f in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {f}"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    p.font.name = 'Calibri'
    p.space_after = Pt(10)

# Right side - Tech stack cards
add_txt(slide, "TECH STACK", 6.5, 0.6, 6, 0.4, 13, ACCENT, False)
add_txt(slide, "Technology Stack", 6.5, 1.0, 6, 0.6, 26, DARK1, True)
add_accent_bar(slide, 6.5, 1.6, 0.8)

tech = [
    ("Backend", "Python 3.12 • FastAPI\nSQLAlchemy • Celery • Redis", DARK3),
    ("Frontend", "Next.js 16 • React 19\nTypeScript • Tailwind CSS v4", DARK3),
    ("Database", "PostgreSQL 16 • SQLite\nQdrant Vector DB • Redis Cache", DARK3),
    ("AI / Agents", "OpenAI GPT-4o\nAgents SDK • 5 Agents", DARK3),
    ("DevOps", "Docker • Git • Ruff\nmypy • pytest • Vitest", DARK3),
    ("Notifications", "Twilio (SMS) • SendGrid\nEasyPost • SmartyStreets", DARK3),
]
for i, (title, desc, bg) in enumerate(tech):
    col = i % 2
    row = i // 2
    x = 6.5 + col * 3.3
    y = 2.0 + row * 1.75
    add_card(slide, title, desc, x, y, 3.0, 1.5, bg, ACCENT)

# ========================
# SLIDE 3: System Architecture
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_txt(slide, "✦  ARCHITECTURE", 0.6, 0.4, 6, 0.4, 13, ACCENT, False)
add_txt(slide, "System Architecture", 0.6, 0.8, 6, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.4, 0.8)

arch = [
    ("⚛", "Frontend", "Next.js 16 • React 19\nTailwind CSS v4\nPort 3000", RGBColor(0x0D, 0x25, 0x4A), ACCENT),
    ("⚡", "Backend API", "FastAPI • Uvicorn\nREST + WebSocket\nPort 8000", RGBColor(0x0A, 0x2A, 0x20), GREEN),
    ("🤖", "AI Agents", "OpenAI Agents SDK\n5 Specialized Agents\nGPT-4o Powered", RGBColor(0x1E, 0x0A, 0x3A), ACCENT2),
    ("🗄", "Vector DB", "Qdrant\n4 Collections\n1536-dim Vectors", RGBColor(0x1E, 0x0A, 0x3A), ACCENT2),
    ("🗃", "PostgreSQL 16", "Relational DB\nOrders, Shipments\nFulfillment Centers", RGBColor(0x0A, 0x2A, 0x20), GREEN),
    ("📦", "Redis 7", "Message Broker\nCelery Backend\nCache Layer", RGBColor(0x0A, 0x2A, 0x20), GREEN),
    ("🔄", "Celery Workers", "Async Tasks\nMonitor Cycles\nBackground Jobs", RGBColor(0x2A, 0x1A, 0x0A), ORANGE),
    ("🌐", "External APIs", "Twilio • SendGrid\nOpenAI • EasyPost\nSmartyStreets", RGBColor(0x2A, 0x0A, 0x0A), RED),
]
for i, (icon, title, desc, bg, color) in enumerate(arch):
    x = 0.5 + (i % 4) * 3.2
    y = 1.8 + (i // 4) * 1.9
    add_card(slide, f"{icon}  {title}", desc, x, y, 2.9, 1.6, bg, color)

add_txt(slide, "Data Flow:  User → Frontend (Next.js) → API (FastAPI) → AI Agents → DB/External Services → Response", 0.5, 5.8, 12.5, 0.5, 13, TEXT2, False, PP_ALIGN.CENTER)
add_txt(slide, "Deployment: Docker Containers (PostgreSQL + Redis + API + Celery Workers + Celery Beat)", 0.5, 6.3, 12.5, 0.5, 13, TEXT2, False, PP_ALIGN.CENTER)

# ========================
# SLIDE 4: AI Agents
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, DARK2, Inches(0), Inches(0), Inches(13.333), Inches(1.4))

add_txt(slide, "✦  AGENTS", 0.6, 0.35, 6, 0.4, 13, ACCENT)
add_txt(slide, "AI Agents Architecture", 0.6, 0.7, 10, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.2, 0.8)

agents = [
    ("Fulfillment Orchestrator", "Central coordinator that runs the full monitor cycle: checks delays, evaluates reroutes, predicts failures, and optimizes costs across all agents.", ACCENT),
    ("Routing Agent", "Selects optimal fulfillment center & carrier rate for new orders based on cost, location, and SLA requirements.", GREEN),
    ("Monitor Agent", "Queries active shipments every N seconds, detects delays and overdue status, and triggers alerts.", ORANGE),
    ("Rerouting Agent", "Evaluates alternative carriers when delays occur, checks feasibility, and executes automated reroutes.", RED),
    ("Communication Agent", "Sends delay alerts via email (SendGrid) and SMS (Twilio) with customer-preferred channels.", ACCENT),
    ("Prediction Agent", "Predicts failure probability for shipments based on delay history, carrier statistics, and risk factors using vector similarity.", GREEN),
    ("Cost Optimizer", "Analyzes shipping costs per monitor cycle, identifies trends, and generates cost-reduction recommendations.", ORANGE),
]
for i, (name, desc, color) in enumerate(agents):
    y = 1.7 + i * 0.8
    add_rect(slide, color, Inches(0.5), Inches(y), Inches(0.06), Inches(0.55))
    add_txt(slide, name, 0.8, y, 3.5, 0.35, 16, WHITE, True)
    add_txt(slide, desc, 4.5, y, 8.2, 0.55, 12, TEXT2)

# ========================
# SLIDE 5: Guardrails
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_txt(slide, "✦  SAFETY", 0.6, 0.5, 6, 0.4, 13, ACCENT)
add_txt(slide, "Business Guardrails", 0.6, 0.9, 8, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.5, 0.8)
add_txt(slide, "Safety rules enforced during every monitor cycle to ensure business compliance", 0.6, 1.7, 11, 0.4, 14, TEXT2)

guardrails = [
    ("🛡  SLA Compliance", "Flags shipments exceeding\ncritical SLA hours thresholds"),
    ("💰  Cost Cap", "Prevents reroutes exceeding\nmax allowed cost increase %"),
    ("🔕  Notification Frequency", "Limits notifications per order\nto avoid spam (max 4)"),
    ("⚠  Failed Delivery Threshold", "Flags shipments exceeding\noverdue delivery threshold"),
    ("🔀  Carrier Diversity", "Prevents switching to same\nmonopoly carriers (>70% vol)"),
    ("📍  Address Validation", "Validates address format,\nZIP code, and street indicators"),
]
for i, (title, desc) in enumerate(guardrails):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.2
    y = 2.3 + row * 2.3
    add_card(slide, title, desc, x, y, 3.8, 1.8, DARK3, ACCENT)

add_txt(slide, "⚡ All guardrails run asynchronously during each Celery Beat monitor cycle (every 15 minutes)", 0.6, 6.7, 12, 0.4, 12, TEXT3, False, PP_ALIGN.CENTER)

# ========================
# SLIDE 6: Frontend
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, DARK2, Inches(0), Inches(0), Inches(13.333), Inches(1.4))

add_txt(slide, "✦  UI", 0.6, 0.35, 6, 0.4, 13, ACCENT)
add_txt(slide, "Frontend Highlights", 0.6, 0.7, 10, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.2, 0.8)

features = [
    ("🏠 Landing Page", "3D Warehouse Scene (Three.js)\nHero animations\nSmooth scrolling (Lenis)\nParallax, Tilt cards"),
    ("📊 Dashboard", "Real-time order tracking\nInventory zones panel\nAnalytics & metrics\nAgent activity logs"),
    ("🤖 AI Assistant", "Chat (GPT-4o-mini)\nStreaming responses\nSmart suggestions\nContext-aware replies"),
    ("🎨 UI/UX", "Radix UI components\nTailwind CSS v4\nFramer Motion\nGSAP timelines"),
    ("🔐 Auth & Forms", "NextAuth.js\nOAuth (Google + FB)\nReact Hook Form + Zod\nJWT sessions"),
    ("✨ Visual Effects", "Noise overlay • Cursor\nMagnetic buttons\nRotating borders\nText reveals"),
    ("📄 Pages", "/dashboard • /orders\n/agents • /analytics\n/workflows • /monitoring\n/team"),
    ("🧪 Testing", "Vitest • RTL\nJest DOM matchers\nUser event simulation\njsdom environment"),
]
for i, (title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.2
    y = 1.7 + row * 2.8
    c = add_card(slide, title, desc, x, y, 2.9, 2.4, DARK3, ACCENT)
    if row == 1:
        for shape in slide.shapes:
            pass

# ========================
# SLIDE 7: Tech Stack
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_txt(slide, "✦  TECH", 0.6, 0.4, 6, 0.4, 13, ACCENT)
add_txt(slide, "Technology Stack", 0.6, 0.8, 8, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.4, 0.8)

stacks = [
    ("Languages", "Python 3.12, TypeScript 5, Node.js 20+, pnpm 9+"),
    ("Backend", "FastAPI, Uvicorn, Pydantic, SQLAlchemy 2.0 (async)"),
    ("Frontend", "Next.js 16 (App Router), React 19, Tailwind CSS v4"),
    ("Database & Cache", "PostgreSQL 16, SQLite (dev), Redis 7, Qdrant Vector DB"),
    ("AI / LLM", "OpenAI GPT-4o, OpenAI Agents SDK, text-embedding-3-small"),
    ("Task Queue", "Celery, Celery Beat, Redis (broker + backend)"),
    ("Auth", "NextAuth.js, JWT (python-jose), OAuth2"),
    ("UI Components", "Radix UI, lucide-react, class-variance-authority"),
    ("Animations", "Framer Motion 12, GSAP 3, Three.js, Lenis"),
    ("Forms & Validation", "React Hook Form, Zod v4, @hookform/resolvers"),
    ("Notifications", "Twilio (SMS), SendGrid (Email)"),
    ("External APIs", "OpenAI, EasyPost (Shipping), SmartyStreets (Address)"),
    ("DevOps", "Docker, docker-compose, Git, ESLint, Ruff, mypy"),
    ("Testing", "pytest-asyncio, Vitest, React Testing Library"),
]
for i, (cat, val) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.7 + row * 0.38
    add_txt(slide, f"▸ {cat}", x, y, 2.2, 0.3, 11, ACCENT, True)
    add_txt(slide, val, x + 2.3, y, 3.8, 0.3, 11, TEXT2)

# ========================
# SLIDE 8: Database
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, DARK2, Inches(0), Inches(0), Inches(13.333), Inches(1.4))

add_txt(slide, "✦  DATA", 0.6, 0.35, 6, 0.4, 13, ACCENT)
add_txt(slide, "Database & Vector Store", 0.6, 0.7, 10, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.2, 0.8)

add_txt(slide, "Relational Database", 0.5, 1.6, 5, 0.4, 16, WHITE, True)
db_rel = [
    ("📋 Orders", "Order ID, Customer, Items\nStatus, Total, Timeline\nAgent Assignments"),
    ("📦 Shipments", "Shipment ID, Carrier\nTracking, Status\nDelays, Reroutes"),
    ("🏭 Fulfillment Centers", "Location, Capacity\nActive Hours\nCarrier Partners"),
    ("💰 Carrier Rates", "Rate Tables\nService Levels\nCost per Zone"),
]
for i, (title, desc) in enumerate(db_rel):
    x = 0.5 + i * 3.2
    add_card(slide, title, desc, x, 2.1, 2.9, 1.7, RGBColor(0x0D, 0x25, 0x4A), ACCENT)

add_txt(slide, "Vector Database (Qdrant) — 1536-dimension embeddings", 0.5, 4.1, 8, 0.4, 16, WHITE, True)
db_vec = [
    ("🚚 shipment_events", "Vector embeddings of\nshipment tracking events\nfor similarity search"),
    ("📦 product_catalog", "Product descriptions &\nfeatures as vectors\nfor semantic search"),
    ("👤 customer_order_history", "Customer patterns &\npreferences for\npersonalized routing"),
    ("🧠 agent_decisions", "Agent reasoning traces\nfor audit & continuous\nimprovement"),
]
for i, (title, desc) in enumerate(db_vec):
    x = 0.5 + i * 3.2
    add_card(slide, title, desc, x, 4.6, 2.9, 1.7, RGBColor(0x1E, 0x0A, 0x3A), ACCENT2)

add_txt(slide, "⚡ Mock fallback available — Qdrant not required for local development", 0.5, 6.6, 12, 0.3, 11, TEXT3)

# ========================
# SLIDE 9: Workflow
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_txt(slide, "✦  WORKFLOW", 0.6, 0.4, 6, 0.4, 13, ACCENT)
add_txt(slide, "Order Fulfillment Workflow", 0.6, 0.8, 10, 0.6, 28, WHITE, True)
add_accent_bar(slide, 0.6, 1.4, 0.8)

steps = [
    ("1", "Order Placed", "Customer places order\nvia frontend UI", ACCENT),
    ("2", "Routing Agent", "Selects best fulfillment\ncenter & carrier rate", GREEN),
    ("3", "Order Processing", "Order assigned to FC\nInventory reserved", ORANGE),
    ("4", "Fulfillment", "Items picked & packed\nShipping label generated", RED),
    ("5", "In Transit", "Real-time tracking\nMonitor agent active", ACCENT),
    ("6", "Delay Detection", "Delay predicted via\nvector similarity search", GREEN),
    ("7", "Rerouting", "Alternative carrier\nselected & executed", ORANGE),
    ("8", "Communication", "SMS/Email alerts sent\nvia Twilio & SendGrid", RED),
    ("9", "Delivery", "Order delivered\nCustomer notified", GREEN),
]
for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.7 + row * 1.9

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.5), Inches(y), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    add_txt(slide, title, x, y + 0.75, 3.8, 0.35, 16, color, True, PP_ALIGN.CENTER)
    add_txt(slide, desc, x, y + 1.1, 3.8, 0.5, 12, TEXT2, False, PP_ALIGN.CENTER)

# ========================
# SLIDE 10: Thank You
# ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, RGBColor(0x05, 0x09, 0x12))
add_rect(slide, ACCENT, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))

add_txt(slide, "✦  FINAL", 1.5, 1.5, 10, 0.5, 13, ACCENT, False, PP_ALIGN.CENTER)
add_txt(slide, "THANK YOU", 1.5, 2.2, 10, 1.2, 52, WHITE, True, PP_ALIGN.CENTER)
add_accent_bar(slide, 5.5, 3.5, 2.2)
add_txt(slide, "Warehouse OS — Multi-Agent Order Fulfillment System", 1.5, 3.9, 10, 0.6, 22, TEXT2, False, PP_ALIGN.CENTER)

add_rect(slide, DARK3, Inches(4.5), Inches(4.8), Inches(4.3), Inches(0.5), 0.15)
add_txt(slide, "📂  Final Project Presentation", 4.6, 4.85, 4.1, 0.4, 14, ACCENT, False, PP_ALIGN.CENTER)

add_txt(slide, "Questions?", 1.5, 5.8, 10, 1, 32, TEXT3, False, PP_ALIGN.CENTER)

# Save
output_path = "C:\\Users\\AC\\Desktop\\final project\\Warehouse_OS_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
